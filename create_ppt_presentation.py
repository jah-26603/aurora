# -*- coding: utf-8 -*-
"""
Created on Wed Dec  4 12:00:51 2024

@author: JDawg
"""
from pptx import Presentation
from pptx.util import Inches
import os
import numpy as np
import glob
from tqdm import tqdm

def create_ppt_with_images(image_paths, prs):
    """
    Add slides with a 3x3 grid of images to the given PowerPoint presentation.
    """
    n_rows, n_cols = 3, 3
    img_width, img_height = Inches(10/3), Inches(15/6)
    x_spacing, y_spacing = Inches(0), Inches(0)

    for i in range(0, len(image_paths), n_rows * n_cols):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
        grid_images = image_paths[i:i + n_rows * n_cols]

        for idx, img_path in enumerate(grid_images):
            row, col = divmod(idx, n_cols)
            x = col * (img_width + x_spacing)
            y = row * (img_height + y_spacing)
            slide.shapes.add_picture(img_path, x, y, width=img_width, height=img_height)
            
import tempfile
import shutil

def combine_presentations(input_ppts, output_ppt):
    """
    Combine multiple PowerPoint presentations into a single file.
    
    Args:
        input_ppts (list of str): List of input PowerPoint file paths to combine.
        output_ppt (str): Path to save the combined PowerPoint file.
    """
    # Create a new, empty presentation
    combined_prs = Presentation()

    for ppt_path in input_ppts:
        current_prs = Presentation(ppt_path)

        for slide in current_prs.slides:
            # Append slide to the combined presentation
            slide_layout = combined_prs.slide_layouts[5]  # Use blank slide layout
            new_slide = combined_prs.slides.add_slide(slide_layout)
            
            # Copy all elements from the current slide
            for shape in slide.shapes:
                if shape.shape_type == 13:  # If it's a picture
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    
                    # Save the image blob to a temporary file
                    with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
                        tmp_file.write(shape.image.blob)
                        tmp_file.close()  # Close the temporary file
                        
                        # Add the image from the temporary file
                        new_slide.shapes.add_picture(tmp_file.name, left, top, width, height)
                        
                        # Remove the temporary file after using it
                        os.remove(tmp_file.name)
                
                elif shape.has_text_frame:  # If it has text
                    textbox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                    textbox.text = shape.text

    # Save the combined presentation
    combined_prs.save(output_ppt)
    print(f"Combined presentation saved as {output_ppt}")


if __name__ == "__main__":

    import re  # Added for regex-based number extraction
    
    pfp = r'D:\gold_level1c_2020_every_7th'
    species = ['1356', '1493', 'LBH']
    results = ['raw_north', 'difference', 'results']
    
    current_dir = pfp
    valid_folders = [os.path.join(current_dir, folder) for folder in os.listdir(current_dir) 
                     if os.path.isdir(os.path.join(current_dir, folder)) and folder.isdigit()]
    
    all_folders = np.array([[[os.path.join(parent, 'graphics', spec, f'{f}') 
                              for f in results] 
                             for spec in species] 
                            for parent in valid_folders])
    
    all_folders = all_folders.reshape(all_folders.shape[0], 9)
    num_days = all_folders.shape[0]
    
    # Function to extract the last numeric part before ".png"
    def extract_sort_key(fp):
        filename = os.path.basename(fp)  # Extract filename from full path
        match = re.search(r'(\d{2}_\d{2})', filename)  # Find the pattern NN_NN (like 08_11)
        if match:
            num_part = match.group(1)  # Extract matched number
            num_value = tuple(map(int, num_part.split('_')))  # Convert "08_11" → (8, 11) for sorting
            return num_value
        return (float('inf'), float('inf'))  # Default high value for non-matching files
    
    # Sort the filepaths numerically based on extracted values
    all_filepaths = [[sorted(glob.glob(f'{parent}/*.png'), key=extract_sort_key)] 
                     for parent in all_folders.flatten()]
    
    all_filepaths = [fp for sublist in all_filepaths for fp in sublist]
    batch_size = 12
    for batch_start in tqdm(range(0, num_days - 1, batch_size)):
        prs = Presentation()
        for i in range(batch_start, min(batch_start + batch_size, num_days - 1)):
            try:
                c_fp = np.array(all_filepaths[int(9 * i):int(9 * (i + 1))])
            except ValueError:
                print(f"Skipping day {i} due to mismatched data.")
                continue
            if c_fp.shape[1] == 0:
                continue
            for scan in range(c_fp.shape[1]):
                images = c_fp[:, scan]
                create_ppt_with_images(images, prs)
        nn = batch_start // batch_size
        mm = str(nn).zfill(3)
        output_ppt = f"2020_kp_actives{mm}.pptx"
        prs.save(output_ppt)
        print(f"Batch PowerPoint saved as {output_ppt}")



    #     # Example usage
    # input_ppts = [ppt for ppt in glob.glob(f'{os.getcwd()}/*.pptx')]# Add all batch file paths
    # output_ppt = "combined_presentation.pptx"
    # combine_presentations(input_ppts, output_ppt)
    
    
    
    
#import os
# import shutil
# import numpy as np

# # Base directory
# pfp = r'D:\gold_level1c_2020_every_7th'

# # Subfolder structure
# species = ['1356', '1493', 'LBH']
# results = ['raw_north', 'difference', 'results']

# # Folder to copy as a replacement
# SOURCE_FOLDER = r'D:\path\to\replacement_folder'  # Update this path

# # Find all valid numbered directories in the base folder
# valid_folders = [
#     os.path.join(pfp, folder)
#     for folder in os.listdir(pfp)
#     if os.path.isdir(os.path.join(pfp, folder)) and folder.isdigit()
# ]

# # Generate full paths for 'graphics/species/results' folders
# all_folders = np.array([
#     [[os.path.join(parent, 'graphics', spec, f'{f}')
#       for f in results]
#      for spec in species]
#     for parent in valid_folders
# ])

# mask_folders = [os.path.join(parent, 'graphics\masks') for parent in valid_folders]
# # Delete and replace "results" folders
# for day_fp in valid_folders:
#     base = os.path.join(valid_folders,'graphics')
#     mask = os.path.join(base, 'masks')
    
#     for spec in species:
#         species_fp = os.path.join(base, f'{spec}\\results') 
#         if os.path.exists(species_fp):
#             shutil.rmtree(species_fp)

#         # Copy the contents of mask to species_fp
#         shutil.copytree(mask, species_fp)